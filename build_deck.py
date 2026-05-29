#!/usr/bin/env python3
"""
Captello Onboarding Deck builder (python-pptx). 6 slides, all clients (<= 7 cap).
Usage:
    pip3 install python-pptx
    python3 build_deck.py content.json out.pptx

content.json schema:
{
  "client": "Philips",
  "call_type": "returning",            # "returning" or "first" — controls slide-2 framing + progress NOW stage
  "subhead": "Building toward a launch-ready event",
  "recap_points": ["...", "...", "..."], # returning: recap of last meeting; first: client goals
  "package_skus": ["ULC x15", "QuickScan", "Salesforce integration"],  # REQUIRED for every deck
  "session1": {"title": "Set up the foundation", "items": ["...","...","..."]},
  "session2": {"title": "Prepare for the floor", "items": ["...","...","..."]},
  "prep_items": ["Confirm capture form fields", "...", "...", "..."]
}
Logo PNGs must be in ./logos/ relative to CWD.
"""
import json, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

RED=RGBColor(0xFF,0,0); BLACK=RGBColor(0,0,0); WHITE=RGBColor(0xFF,0xFF,0xFF)
GRAY=RGBColor(0x80,0x80,0x80); LGRAY=RGBColor(0xD9,0xD9,0xD9); DGRAY=RGBColor(0x40,0x40,0x40)
FONT="Roboto"; ASPECT=5.415
LOGO_BLACK="logos/captello-horizontal-blacktext-transparent.png"
LOGO_WHITE="logos/captello-horizontal-whitetext-transparent.png"
CONF=("CONFIDENTIALITY NOTICE: This presentation is for the sole use of the intended recipient and "
      "may contain proprietary, confidential and/or trade secret information. Any unauthorized review, "
      "use, disclosure or distribution is prohibited. Thank you.")

def solid(sh,c): sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def sf(r,s,c,b=False): r.font.size=Pt(s); r.font.color.rgb=c; r.font.bold=b; r.font.name=FONT
def tbox(slide,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        for seg in line:
            r=p.add_run(); r.text=seg["t"]; sf(r,seg["s"],seg["c"],seg.get("b",False))
            if seg.get("spc"):
                r._r.get_or_add_rPr().set("spc",str(int(seg["spc"]*100)))
    return tb
def logo(slide,path,x,y,h): slide.shapes.add_picture(path,Inches(x),Inches(y),height=Inches(h))
def arc(slide):
    s=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(7.35),Inches(3.75),Inches(5.6),Inches(5.6)); solid(s,RED)
def bullets(slide,x,y,w,h,items,size,color,space=12):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space)
        pPr=p._p.get_or_add_pPr(); bu=pPr.makeelement(qn('a:buChar'),{'char':'\u2022'}); pPr.append(bu)
        pPr.set('indent','-228600'); pPr.set('marL','228600')
        r=p.add_run(); r.text=it; sf(r,size,color)
    return tb
def headline(slide,segs):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.5),Inches(0.33),Inches(0.07),Inches(0.6)); solid(bar,RED)
    tbox(slide,0.65,0.33,8.85,0.6,[segs],anchor=MSO_ANCHOR.MIDDLE)
def footer(slide):
    ln=slide.shapes.add_connector(2,Inches(0.5),Inches(5.105),Inches(9.5),Inches(5.105)); ln.line.color.rgb=LGRAY; ln.line.width=Pt(0.75)
    logo(slide,LOGO_BLACK,0.5,5.205,0.20)
    tbox(slide,2.1,5.165,7.4,0.4,[[{"t":CONF,"s":6,"c":GRAY}]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
def light(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=WHITE; return s
def dark(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=BLACK; return s

def progress_tracker(slide, now_idx):
    pillX,pillW,pillH,top,gap=6.0,3.3,0.72,1.35,0.42
    stages=["Kicked off","Working sessions","Event-ready"]
    ln=slide.shapes.add_connector(2,Inches(pillX+pillW/2),Inches(top+pillH/2),Inches(pillX+pillW/2),Inches(top+pillH/2+2*(pillH+gap)))
    ln.line.color.rgb=LGRAY; ln.line.width=Pt(1.5)
    for i,t in enumerate(stages):
        y=top+i*(pillH+gap); cur=(i==now_idx)
        box=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(pillX),Inches(y),Inches(pillW),Inches(pillH))
        if cur: solid(box,RED)
        else: box.fill.solid(); box.fill.fore_color.rgb=WHITE; box.line.color.rgb=LGRAY; box.line.width=Pt(1)
        tf=box.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        if cur:
            r=p.add_run(); r.text="NOW  \u00b7  "; sf(r,11,WHITE,True)
        r=p.add_run(); r.text=t; sf(r,16,WHITE if cur else DGRAY,cur)

def package_slide(prs, skus):
    s=light(prs)
    headline(s,[{"t":"Your Package, ","s":28,"c":BLACK,"b":True},{"t":"At a Glance","s":28,"c":RED,"b":True},{"t":".","s":28,"c":BLACK,"b":True}])
    if not skus: skus=["See contract for purchased services"]
    cols=2; cwd=4.25; gx=0.5; x0=0.65; y0=1.6
    n=len(skus); rows=(n+cols-1)//cols
    ch=min(0.85, max(0.5,(3.3-(rows-1)*0.28)/max(rows,1))); gy=0.28
    for i,sku in enumerate(skus):
        r=i//cols; c=i%cols; x=x0+c*(cwd+gx); y=y0+r*(ch+gy)
        card=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(cwd),Inches(ch))
        card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0xF7,0xF7,0xF7); card.line.color.rgb=RGBColor(0xEA,0xEA,0xEA)
        bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(0.06),Inches(ch)); solid(bar,RED)
        tf=card.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.margin_left=Inches(0.25)
        p=tf.paragraphs[0]; rr=p.add_run(); rr.text=sku; sf(rr,16,BLACK,True)
    footer(s)

def build(c,out):
    first=c.get("call_type")=="first"
    prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)

    # 1 COVER
    s=dark(prs); arc(s); logo(s,LOGO_WHITE,0.5,0.45,0.30)
    tbox(s,0.5,2.05,8,0.35,[[{"t":"ONBOARDING   \u00b7   WORKING SESSION","s":13,"c":RED,"b":True,"spc":3}]])
    tbox(s,0.48,2.35,8,1.0,[[{"t":c["client"],"s":54,"c":WHITE,"b":True}]])
    tbox(s,0.5,3.45,7.5,0.45,[[{"t":c.get("subhead","Building toward a launch-ready event"),"s":20,"c":RGBColor(0xE0,0xE0,0xE0)}]])
    tbox(s,0.5,4.95,7,0.35,[[{"t":"Prepared by Captello Onboarding & Enablement","s":11,"c":GRAY}]])

    # 2 RECAP (returning) / GOALS (first) -- REQUIRED
    s=light(prs)
    if first:
        headline(s,[{"t":"Your Goals, ","s":28,"c":BLACK,"b":True},{"t":"Our Plan","s":28,"c":RED,"b":True},{"t":".","s":28,"c":BLACK,"b":True}])
        now_idx=0
        pts=c.get("recap_points",["Capture qualified leads at your next event","Equip staff with a fast, reliable workflow","Sync results into your CRM automatically"])
    else:
        headline(s,[{"t":"Where We ","s":28,"c":BLACK,"b":True},{"t":"Left Off","s":28,"c":RED,"b":True},{"t":".","s":28,"c":BLACK,"b":True}])
        now_idx=1
        pts=c.get("recap_points",["Account live and configured","Team onboarded on core workflow","Two working sessions ahead"])
    bullets(s,0.65,1.55,4.7,2.6,pts,19,BLACK,14)
    progress_tracker(s,now_idx)
    footer(s)

    # 3 PACKAGE -- REQUIRED for ALL clients
    package_slide(prs, c.get("package_skus",[]))

    # 4 AGENDA
    s=light(prs)
    headline(s,[{"t":"Two Sessions to ","s":28,"c":BLACK,"b":True},{"t":"Launch-Ready","s":28,"c":RED,"b":True},{"t":".","s":28,"c":BLACK,"b":True}])
    s1=c.get("session1",{"title":"Set up the foundation","items":["Verify integration","Build capture form","Set qualification questions"]})
    s2=c.get("session2",{"title":"Prepare for the floor","items":["Test scan and badges","Invite and license staff","Confirm follow-up and CRM sync"]})
    colY,colW,colH,gap=1.55,4.25,2.95,0.5; c1x=0.65; c2x=c1x+colW+gap
    def card(x,label,title,items):
        bx=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(colY),Inches(colW),Inches(colH))
        bx.fill.solid(); bx.fill.fore_color.rgb=RGBColor(0xF7,0xF7,0xF7); bx.line.color.rgb=RGBColor(0xEA,0xEA,0xEA)
        top=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(colY),Inches(colW),Inches(0.06)); solid(top,RED)
        tbox(s,x+0.3,colY+0.22,colW-0.6,0.3,[[{"t":label,"s":12,"c":RED,"b":True,"spc":2}]])
        tbox(s,x+0.3,colY+0.5,colW-0.6,0.4,[[{"t":title,"s":18,"c":BLACK,"b":True}]])
        bullets(s,x+0.3,colY+1.05,colW-0.6,colH-1.25,items,15,RGBColor(0x30,0x30,0x30),10)
    card(c1x,"SESSION 1",s1["title"],s1["items"]); card(c2x,"SESSION 2",s2["title"],s2["items"])
    footer(s)

    # 5 PREPARATION
    s=light(prs)
    headline(s,[{"t":"Your Prep Drives a ","s":28,"c":BLACK,"b":True},{"t":"Flawless","s":28,"c":RED,"b":True},{"t":" Event.","s":28,"c":BLACK,"b":True}])
    items=c.get("prep_items",["Confirm capture form fields","Finalize qualification questions","Provide staff list for licenses","Approve follow-up email content"])
    TEXT_H=0.42; ICON=0.34; off=(TEXT_H-ICON)/2; rs=0.66
    groupH=(len(items)-1)*rs+TEXT_H; startY=1.5+(4.95-1.5-groupH)/2
    for i,it in enumerate(items):
        y=startY+i*rs
        ic=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.65),Inches(y+off),Inches(ICON),Inches(ICON)); solid(ic,RED)
        tf=ic.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text="\u2713"; sf(r,15,WHITE,True)
        tbox(s,0.65+ICON+0.18,y,4.6,TEXT_H,[[{"t":it,"s":17,"c":BLACK}]],anchor=MSO_ANCHOR.MIDDLE)
    cx,cy,cwd2,chd=6.0,1.55,3.45,2.9
    cardd=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(cx),Inches(cy),Inches(cwd2),Inches(chd)); solid(cardd,BLACK)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(cx),Inches(cy),Inches(0.07),Inches(chd)); solid(bar,RED)
    tbox(s,cx+0.35,cy+0.5,cwd2-0.6,0.9,[[{"t":"Test before launch","s":24,"c":WHITE,"b":True}]])
    tbox(s,cx+0.35,cy+1.5,cwd2-0.6,1.2,[[{"t":"We validate every setup end-to-end before your event goes live — so day-of runs clean.","s":14,"c":RGBColor(0xDA,0xDA,0xDA)}]])
    footer(s)

    # 6 CLOSE
    s=dark(prs); arc(s); logo(s,LOGO_WHITE,0.5,0.45,0.30)
    tbox(s,0.5,2.25,8.4,1.4,[[{"t":"Let's make your next event your best.","s":38,"c":WHITE,"b":True}]],anchor=MSO_ANCHOR.MIDDLE)
    tbox(s,0.5,3.75,8,0.8,[[{"t":"Your onboarding specialist is one message away.","s":16,"c":RGBColor(0xE0,0xE0,0xE0)}],[{"t":"captello.com","s":16,"c":RED,"b":True}]])

    prs.save(out); print("WROTE",out,"slides:",len(prs.slides._sldIdLst))

if __name__=="__main__":
    c=json.load(open(sys.argv[1])) if len(sys.argv)>1 else {"client":"Sample","call_type":"returning"}
    build(c, sys.argv[2] if len(sys.argv)>2 else "deck.pptx")
